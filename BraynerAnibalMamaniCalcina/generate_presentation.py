import os
import sys
import subprocess

# Auto-instalar python-pptx si no está presente
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("[*] python-pptx no encontrado. Instalando automáticamente...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Configurar diapositivas a formato moderno 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Paleta de colores Premium (Slate Theme con acentos Sky Blue y Orange)
    COLOR_BG = RGBColor(15, 23, 42)          # Slate 900 (Fondo)
    COLOR_CARD = RGBColor(30, 41, 59)        # Slate 800 (Tarjetas/Contenedores)
    COLOR_PRIMARY = RGBColor(56, 189, 248)   # Sky Blue 400 (Títulos/Acentos principales)
    COLOR_TEXT = RGBColor(241, 245, 249)     # Slate 100 (Texto general)
    COLOR_MUTED = RGBColor(148, 163, 184)    # Slate 400 (Subtítulos/Texto secundario)
    COLOR_ACCENT = RGBColor(249, 115, 22)    # Orange 500 (Alertas/Destacados)
    COLOR_GREEN = RGBColor(34, 197, 94)      # Green 500 (Casos de prueba/Aprobado)

    def apply_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG
        
        # Barra superior decorativa
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(13.333), Inches(0.15)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = COLOR_PRIMARY
        accent_bar.line.fill.background()

    def add_title(slide, text):
        txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.font.name = "Arial"

    # ==========================================
    # DIAPOSITIVA 1: Portada (Carátula)
    # ==========================================
    slide_layout = prs.slide_layouts[6] # Blank
    slide1 = prs.slides.add_slide(slide_layout)
    apply_background(slide1)
    
    # Título Principal (YatiqApp)
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.33), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "YatiqApp"
    p1.font.size = Pt(72)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_PRIMARY
    p1.font.name = "Arial"
    
    p2 = tf.add_paragraph()
    p2.text = "Aprendizaje de Aimara y Quechua"
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXT
    p2.font.name = "Arial"
    p2.space_before = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = "Asistente Conversacional Inteligente para la Preservación Lingüística"
    p3.font.size = Pt(18)
    p3.font.color.rgb = COLOR_MUTED
    p3.font.name = "Arial"
    p3.space_before = Pt(15)

    # Rectángulo Decorativo del Autor
    author_box = slide1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.0), Inches(4.8), Inches(7.5), Inches(1.5)
    )
    author_box.fill.solid()
    author_box.fill.fore_color.rgb = COLOR_CARD
    author_box.line.color.rgb = COLOR_PRIMARY
    author_box.line.width = Pt(1.5)
    
    tf_author = author_box.text_frame
    tf_author.word_wrap = True
    tf_author.margin_left = Inches(0.2)
    tf_author.margin_top = Inches(0.15)
    
    pa1 = tf_author.paragraphs[0]
    pa1.text = "Línea de Evaluación: CE02 - Ingeniería de Software"
    pa1.font.size = Pt(16)
    pa1.font.bold = True
    pa1.font.color.rgb = COLOR_PRIMARY
    
    pa2 = tf_author.add_paragraph()
    pa2.text = "Responsable: Brayner Anibal Mamani Calcina"
    pa2.font.size = Pt(16)
    pa2.font.color.rgb = COLOR_TEXT
    pa2.space_before = Pt(5)
    
    pa3 = tf_author.add_paragraph()
    pa3.text = "Universidad Nacional del Altiplano"
    pa3.font.size = Pt(14)
    pa3.font.color.rgb = COLOR_MUTED
    pa3.space_before = Pt(5)

    # ==========================================
    # DIAPOSITIVA 2: El Problema y La Solución
    # ==========================================
    slide2 = prs.slides.add_slide(slide_layout)
    apply_background(slide2)
    add_title(slide2, "El Problema y la Oportunidad de YatiqApp")

    # Contenedor Izquierdo: El Problema
    prob_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.6), Inches(5.6), Inches(5.0))
    prob_box.fill.solid()
    prob_box.fill.fore_color.rgb = COLOR_CARD
    prob_box.line.color.rgb = COLOR_ACCENT
    prob_box.line.width = Pt(2)
    
    tf_prob = prob_box.text_frame
    tf_prob.word_wrap = True
    tf_prob.margin_left = tf_prob.margin_right = Inches(0.3)
    tf_prob.margin_top = Inches(0.3)
    
    p = tf_prob.paragraphs[0]
    p.text = "❌ EL PROBLEMA"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    
    bullets_prob = [
        "Brecha de Inclusión Digital: Más de 4 millones de peruanos y bolivianos que hablan aimara o quechua están marginados del ecosistema digital actual.",
        "Asistentes Inadecuados: Herramientas líderes de mercado (Siri, Alexa, Google Translate) carecen de soporte de voz interactivo y fluido para lenguas andinas.",
        "Conectividad Nula en Zonas Rurales: Las escuelas y comunidades de Puno no cuentan con internet estable, bloqueando soluciones basadas 100% en la nube."
    ]
    for b in bullets_prob:
        bp = tf_prob.add_paragraph()
        bp.text = "• " + b.split(":")[0] + ":"
        bp.font.size = Pt(16)
        bp.font.bold = True
        bp.font.color.rgb = COLOR_TEXT
        bp.space_before = Pt(15)
        
        bp_desc = tf_prob.add_paragraph()
        bp_desc.text = b.split(":")[1].strip()
        bp_desc.font.size = Pt(14)
        bp_desc.font.color.rgb = COLOR_MUTED
        bp_desc.space_before = Pt(2)

    # Contenedor Derecho: La Solución
    sol_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.98), Inches(1.6), Inches(5.6), Inches(5.0))
    sol_box.fill.solid()
    sol_box.fill.fore_color.rgb = COLOR_CARD
    sol_box.line.color.rgb = COLOR_GREEN
    sol_box.line.width = Pt(2)
    
    tf_sol = sol_box.text_frame
    tf_sol.word_wrap = True
    tf_sol.margin_left = tf_sol.margin_right = Inches(0.3)
    tf_sol.margin_top = Inches(0.3)
    
    p = tf_sol.paragraphs[0]
    p.text = "💡 LA PROPUESTA (YatiqApp)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN
    
    bullets_sol = [
        "Inferencia Híbrida Inteligente: Combinación de modelos locales de última generación (ASR, NMT, TTS) y APIs de LLM externas para orquestar la lógica conversacional.",
        "Aceleración por GPU Local: Ejecución de modelos pesados de Machine Learning en hardware GPU (RTX 5060) logrando traducción en tiempo real (<0.8 segundos).",
        "Resiliencia y Modo Offline: Capacidad de realizar traducciones directas e interacciones locales sin requerir conectividad activa a internet."
    ]
    for b in bullets_sol:
        bp = tf_sol.add_paragraph()
        bp.text = "• " + b.split(":")[0] + ":"
        bp.font.size = Pt(16)
        bp.font.bold = True
        bp.font.color.rgb = COLOR_TEXT
        bp.space_before = Pt(15)
        
        bp_desc = tf_sol.add_paragraph()
        bp_desc.text = b.split(":")[1].strip()
        bp_desc.font.size = Pt(14)
        bp_desc.font.color.rgb = COLOR_MUTED
        bp_desc.space_before = Pt(2)

    # ==========================================
    # DIAPOSITIVA 3: Especificación de Requerimientos (SRS)
    # ==========================================
    slide3 = prs.slides.add_slide(slide_layout)
    apply_background(slide3)
    add_title(slide3, "Requerimientos de Software (CE021)")

    # Card Izquierda: Requerimientos Funcionales (RF)
    rf_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.6), Inches(5.6), Inches(5.0))
    rf_card.fill.solid()
    rf_card.fill.fore_color.rgb = COLOR_CARD
    rf_card.line.color.rgb = COLOR_PRIMARY
    rf_card.line.width = Pt(1.5)
    
    tf_rf = rf_card.text_frame
    tf_rf.word_wrap = True
    tf_rf.margin_left = tf_rf.margin_right = Inches(0.3)
    tf_rf.margin_top = Inches(0.3)
    
    p = tf_rf.paragraphs[0]
    p.text = "📋 Requerimientos Funcionales Clave"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    rfs = [
        "RF-01 Transcripción de Voz (ASR): Conversión de audio del usuario a texto local con Whisper.",
        "RF-02 Traducción Bidireccional: Traducción Español ↔ Aimara/Quechua mediante NLLB-200.",
        "RF-03 Síntesis de Voz (TTS): Conversión de la traducción a audio WAV nativo con Meta MMS.",
        "RF-04 Asistente Inteligente Híbrido: Respuestas conversacionales fluidas utilizando Gemini API.",
        "RF-08 Lecciones e Interactividad (LMS): Módulos educativos estructurados para el auto-aprendizaje."
    ]
    for r in rfs:
        rp = tf_rf.add_paragraph()
        rp.text = r.split(":")[0]
        rp.font.size = Pt(15)
        rp.font.bold = True
        rp.font.color.rgb = COLOR_TEXT
        rp.space_before = Pt(14)
        
        rp_desc = tf_rf.add_paragraph()
        rp_desc.text = r.split(":")[1].strip()
        rp_desc.font.size = Pt(13)
        rp_desc.font.color.rgb = COLOR_MUTED
        rp_desc.space_before = Pt(2)

    # Card Derecha: Requerimientos No Funcionales (RNF)
    rnf_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.98), Inches(1.6), Inches(5.6), Inches(5.0))
    rnf_card.fill.solid()
    rnf_card.fill.fore_color.rgb = COLOR_CARD
    rnf_card.line.color.rgb = COLOR_PRIMARY
    rnf_card.line.width = Pt(1.5)
    
    tf_rnf = rnf_card.text_frame
    tf_rnf.word_wrap = True
    tf_rnf.margin_left = tf_rnf.margin_right = Inches(0.3)
    tf_rnf.margin_top = Inches(0.3)
    
    p = tf_rnf.paragraphs[0]
    p.text = "⚡ Requerimientos No Funcionales"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    rnfs = [
        "RNF-01 Latencia y Velocidad: Pipeline inferencial completo en menos de 1.0s con aceleración CUDA.",
        "RNF-02 Backend Asíncrono: Servidor de alto rendimiento FastAPI para concurrencia multi-cliente.",
        "RNF-03 Portabilidad: Aplicación multiplataforma en Flutter para Android e iOS.",
        "RNF-05 Operación Offline: Lecciones y traducciones básicas almacenadas localmente en SQLite.",
        "RNF-06 Tiempo de Síntesis: Generación de ondas WAV por MMS TTS en menos de 400ms."
    ]
    for r in rnfs:
        rp = tf_rnf.add_paragraph()
        rp.text = r.split(":")[0]
        rp.font.size = Pt(15)
        rp.font.bold = True
        rp.font.color.rgb = COLOR_TEXT
        rp.space_before = Pt(14)
        
        rp_desc = tf_rnf.add_paragraph()
        rp_desc.text = r.split(":")[1].strip()
        rp_desc.font.size = Pt(13)
        rp_desc.font.color.rgb = COLOR_MUTED
        rp_desc.space_before = Pt(2)

    # ==========================================
    # DIAPOSITIVA 4: Arquitectura del Sistema (3-Tier)
    # ==========================================
    slide4 = prs.slides.add_slide(slide_layout)
    apply_background(slide4)
    add_title(slide4, "Arquitectura de Integración de Componentes (CE023)")

    # Dibujar Cajas de Arquitectura (Flutter -> Laravel Proxy -> FastAPI GPU)
    boxes = [
        ("📱 Frontend Móvil / Web", "Flutter (Dart)\nLaravel Blade Views", Inches(0.75), COLOR_PRIMARY),
        ("🛡️ Proxy / Web Server", "Laravel 11 Gateway\nSeguridad y Cookies", Inches(4.9), COLOR_ACCENT),
        ("⚡ Backend Inferencial", "FastAPI (Python)\nUvicorn Asíncrono", Inches(9.05), COLOR_GREEN)
    ]
    
    for title, subtitle, left, color in boxes:
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(3.5), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = color
        card.line.width = Pt(2)
        
        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_right = Inches(0.2)
        tf_c.margin_top = Inches(0.3)
        
        p = tf_c.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf_c.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(15)
        p2.font.color.rgb = COLOR_TEXT
        p2.space_before = Pt(12)
        p2.alignment = PP_ALIGN.CENTER

    # Flechas de conexión
    arrow1 = slide4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.35), Inches(2.7), Inches(0.45), Inches(0.4))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = COLOR_MUTED
    arrow1.line.fill.background()
    
    arrow2 = slide4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.5), Inches(2.7), Inches(0.45), Inches(0.4))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = COLOR_MUTED
    arrow2.line.fill.background()

    # Contenedor de Modelos de IA cargados en GPU (FastAPI)
    gpu_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(4.5), Inches(11.83), Inches(2.1))
    gpu_box.fill.solid()
    gpu_box.fill.fore_color.rgb = COLOR_CARD
    gpu_box.line.color.rgb = COLOR_PRIMARY
    gpu_box.line.width = Pt(1.5)
    
    tf_g = gpu_box.text_frame
    tf_g.word_wrap = True
    tf_g.margin_left = tf_g.margin_right = Inches(0.4)
    tf_g.margin_top = Inches(0.25)
    
    p = tf_g.paragraphs[0]
    p.text = "🧠 Pipeline de Inferencia en Cascada (FastAPI + GPU CUDA)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    p2 = tf_g.add_paragraph()
    p2.text = "1. Reconocimiento de Voz (ASR): OpenAI Whisper Large V3 Turbo transcribe audio WAV a texto."
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_TEXT
    p2.space_before = Pt(10)
    
    p3 = tf_g.add_paragraph()
    p3.text = "2. Traducción Neuronal (NMT): Meta NLLB-200 + Adaptadores LoRA traduce a Aimara/Quechua."
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_TEXT
    p3.space_before = Pt(5)
    
    p4 = tf_g.add_paragraph()
    p4.text = "3. Síntesis de Voz (TTS): Meta MMS (Massively Multilingual Speech) genera archivo de audio en lengua originaria."
    p4.font.size = Pt(14)
    p4.font.color.rgb = COLOR_TEXT
    p4.space_before = Pt(5)

    # ==========================================
    # DIAPOSITIVA 5: Plataforma de Datos (SQLite)
    # ==========================================
    slide5 = prs.slides.add_slide(slide_layout)
    apply_background(slide5)
    add_title(slide5, "Plataforma de Datos y Persistencia (CE022)")

    # Justificación de SQLite
    just_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.6), Inches(5.6), Inches(5.0))
    just_box.fill.solid()
    just_box.fill.fore_color.rgb = COLOR_CARD
    just_box.line.color.rgb = COLOR_PRIMARY
    just_box.line.width = Pt(1.5)
    
    tf_j = just_box.text_frame
    tf_j.word_wrap = True
    tf_j.margin_left = tf_j.margin_right = Inches(0.3)
    tf_j.margin_top = Inches(0.3)
    
    p = tf_j.paragraphs[0]
    p.text = "⚙️ Base de Datos Embebida (SQLite 3)"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    bullet_j = [
        "Cero Administración de Servidores: SQLite elimina la necesidad de procesos en segundo plano independientes, facilitando la portabilidad mediante un único archivo físico: database.sqlite.",
        "Garantías ACID: Asegura atomicidad y consistencia de las transacciones, ideal para la gestión local de colas e historial.",
        "Bajo Consumo de RAM: Libera la memoria del servidor para dedicarla al almacenamiento de los modelos deep learning."
    ]
    for b in bullet_j:
        bp = tf_j.add_paragraph()
        bp.text = "• " + b.split(":")[0] + ":"
        bp.font.size = Pt(15)
        bp.font.bold = True
        bp.font.color.rgb = COLOR_TEXT
        bp.space_before = Pt(15)
        
        bp_desc = tf_j.add_paragraph()
        bp_desc.text = b.split(":")[1].strip()
        bp_desc.font.size = Pt(13)
        bp_desc.font.color.rgb = COLOR_MUTED
        bp_desc.space_before = Pt(2)

    # Gestión de Colas y Hilos
    queue_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.98), Inches(1.6), Inches(5.6), Inches(5.0))
    queue_box.fill.solid()
    queue_box.fill.fore_color.rgb = COLOR_CARD
    queue_box.line.color.rgb = COLOR_PRIMARY
    queue_box.line.width = Pt(1.5)
    
    tf_q = queue_box.text_frame
    tf_q.word_wrap = True
    tf_q.margin_left = tf_q.margin_right = Inches(0.3)
    tf_q.margin_top = Inches(0.3)
    
    p = tf_q.paragraphs[0]
    p.text = "🔄 Gestión de Tareas Asíncronas y VRAM"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    bullet_q = [
        "Cola Persistente (jobs): Laravel utiliza SQLite para encolar el re-entrenamiento del traductor (fine-tuning), evitando caídas o cuelgues del proceso web principal.",
        "Mecanismo Mutex (cache_locks): Semáforo atómico en base de datos que restringe el acceso concurrente a la GPU CUDA, previniendo excepciones de desbordamiento de VRAM.",
        "Historial y Caché: Persistencia local de traducciones frecuentes para acelerar tiempos de respuesta."
    ]
    for b in bullet_q:
        bp = tf_q.add_paragraph()
        bp.text = "• " + b.split(":")[0] + ":"
        bp.font.size = Pt(15)
        bp.font.bold = True
        bp.font.color.rgb = COLOR_TEXT
        bp.space_before = Pt(15)
        
        bp_desc = tf_q.add_paragraph()
        bp_desc.text = b.split(":")[1].strip()
        bp_desc.font.size = Pt(13)
        bp_desc.font.color.rgb = COLOR_MUTED
        bp_desc.space_before = Pt(2)

    # ==========================================
    # DIAPOSITIVA 6: Calidad y Pruebas
    # ==========================================
    slide6 = prs.slides.add_slide(slide_layout)
    apply_background(slide6)
    add_title(slide6, "Aseguramiento de Calidad y Seguridad (CE024)")

    # Matriz de Pruebas
    test_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.6), Inches(5.6), Inches(5.0))
    test_box.fill.solid()
    test_box.fill.fore_color.rgb = COLOR_CARD
    test_box.line.color.rgb = COLOR_GREEN
    test_box.line.width = Pt(1.5)
    
    tf_t = test_box.text_frame
    tf_t.word_wrap = True
    tf_t.margin_left = tf_t.margin_right = Inches(0.3)
    tf_t.margin_top = Inches(0.3)
    
    p = tf_t.paragraphs[0]
    p.text = "🧪 Pruebas Funcionales (100% OK)"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN
    
    tests = [
        "TC-01 Traducción de Texto: Validado Español ↔ Aimara/Quechua.",
        "TC-02 Traducción de Voz: Whisper ASR y MMS TTS integrados.",
        "TC-03 Model Arena: Comparador de adaptadores LoRA y base.",
        "TC-05 Gráfica de Convergencia: Monitoreo visual de la curva de pérdida de entrenamiento.",
        "TC-07 Historial Móvil: Persistencia local en base de datos SQLite."
    ]
    for t in tests:
        tp = tf_t.add_paragraph()
        tp.text = "✓ " + t.split(":")[0]
        tp.font.size = Pt(15)
        tp.font.bold = True
        tp.font.color.rgb = COLOR_TEXT
        tp.space_before = Pt(14)
        
        tp_desc = tf_t.add_paragraph()
        tp_desc.text = t.split(":")[1].strip()
        tp_desc.font.size = Pt(13)
        tp_desc.font.color.rgb = COLOR_MUTED
        tp_desc.space_before = Pt(2)

    # Seguridad y Tolerancia a fallos
    sec_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.98), Inches(1.6), Inches(5.6), Inches(5.0))
    sec_box.fill.solid()
    sec_box.fill.fore_color.rgb = COLOR_CARD
    sec_box.line.color.rgb = COLOR_PRIMARY
    sec_box.line.width = Pt(1.5)
    
    tf_s = sec_box.text_frame
    tf_s.word_wrap = True
    tf_s.margin_left = tf_s.margin_right = Inches(0.3)
    tf_s.margin_top = Inches(0.3)
    
    p = tf_s.paragraphs[0]
    p.text = "🛡️ Seguridad y Resiliencia"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    bullet_s = [
        "Tolerancia a Ruido y Formatos: Conversión automática de canales de audio a mono 16kHz y float32; suspensión inteligente en audios vacíos.",
        "Limpieza Atómica de Datos: Eliminación física inmediata de los archivos de voz (.wav) temporales tras la inferencia para proteger la privacidad del usuario.",
        "Prevención SQL/CSRF: Sentencias preparadas nativas del ORM Eloquent y tokens de seguridad CSRF inyectados en Laravel web."
    ]
    for b in bullet_s:
        bp = tf_s.add_paragraph()
        bp.text = "• " + b.split(":")[0] + ":"
        bp.font.size = Pt(15)
        bp.font.bold = True
        bp.font.color.rgb = COLOR_TEXT
        bp.space_before = Pt(15)
        
        bp_desc = tf_s.add_paragraph()
        bp_desc.text = b.split(":")[1].strip()
        bp_desc.font.size = Pt(13)
        bp_desc.font.color.rgb = COLOR_MUTED
        bp_desc.space_before = Pt(2)

    # ==========================================
    # DIAPOSITIVA 7: Demo Path y Sustentación
    # ==========================================
    slide7 = prs.slides.add_slide(slide_layout)
    apply_background(slide7)
    add_title(slide7, "Guía de Demostración en Vivo (CE0217)")

    steps = [
        ("1. Inicio Consola", "Ejecutar servidores localmente (FastAPI en puerto 8000, Laravel en 8080).", Inches(0.75)),
        ("2. Traducción Voz", "Hablar al micrófono; transcripción Whisper y reproducción por MMS TTS en <0.8s.", Inches(3.2)),
        ("3. Model Arena", "Acceder a `/compare` para contrastar métricas BLEU y ChrF++ de adaptadores LoRA.", Inches(5.65)),
        ("4. Loss Curve", "Revisar gráficas de convergencia históricas en tiempo real en `/reports`.", Inches(8.1)),
        ("5. Demo Móvil", "Proyectar Flutter en Android físico; probar traducción por voz offline sin internet.", Inches(10.55))
    ]
    
    for title, desc, left in steps:
        step_card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.2), Inches(2.1), Inches(4.0))
        step_card.fill.solid()
        step_card.fill.fore_color.rgb = COLOR_CARD
        step_card.line.color.rgb = COLOR_PRIMARY
        step_card.line.width = Pt(1.5)
        
        tf_st = step_card.text_frame
        tf_st.word_wrap = True
        tf_st.margin_left = tf_st.margin_right = Inches(0.15)
        tf_st.margin_top = Inches(0.2)
        
        p = tf_st.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf_st.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = COLOR_TEXT
        p2.space_before = Pt(15)
        p2.alignment = PP_ALIGN.CENTER

    # ==========================================
    # DIAPOSITIVA 8: Roadmap y Futuro
    # ==========================================
    slide8 = prs.slides.add_slide(slide_layout)
    apply_background(slide8)
    add_title(slide8, "Roadmap de Evolución Tecnológica")

    phases = [
        ("🚀 Fase 1: Migración a la Nube (AWS)", "Despliegue de la API inferencial en FastAPI sobre instancias GPU AWS EC2 (g5.xlarge NVIDIA A10G) para soportar múltiples clientes en paralelo con alta escalabilidad.", Inches(1.8)),
        ("🌱 Fase 2: Expansión Lingüística", "Inclusión de nuevas lenguas nativas peruanas (Shipibo-Conibo y Asháninka) mediante adaptadores LoRA ligeros (~14MB) listos para intercambio caliente en caliente (hot-swapping).", Inches(3.4)),
        ("🧠 Fase 3: LLM Conversacional Local", "Entrenamiento e integración de modelos de lenguaje locales (Llama-3-8B cuantizado en 4-bit) para correr sin conexión a internet y reemplazar las llamadas externas de Gemini/OpenAI.", Inches(5.0))
    ]
    
    for title, desc, top in phases:
        ph_card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), top, Inches(11.83), Inches(1.3))
        ph_card.fill.solid()
        ph_card.fill.fore_color.rgb = COLOR_CARD
        ph_card.line.color.rgb = COLOR_PRIMARY
        ph_card.line.width = Pt(1.5)
        
        tf_ph = ph_card.text_frame
        tf_ph.word_wrap = True
        tf_ph.margin_left = tf_ph.margin_right = Inches(0.3)
        tf_ph.margin_top = Inches(0.2)
        
        p = tf_ph.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        
        p2 = tf_ph.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = COLOR_TEXT
        p2.space_before = Pt(5)

    # Guardar Presentación
    output_filename = "YatiqApp_Sustentacion.pptx"
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), output_filename)
    prs.save(output_path)
    print(f"[+] Presentación PowerPoint guardada exitosamente en:\n    {output_path}")

if __name__ == "__main__":
    create_presentation()
